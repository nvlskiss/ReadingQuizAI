import os
import importlib

try:
    import pymupdf
except ImportError:
    try:
        import fitz as pymupdf
    except ImportError:
        pymupdf = None

try:
    import docx
except ImportError:
    docx = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


class ExtractText:
    def __init__(self, path):
        self.path = path or ""
        self.text_format = ""

    def convert(self) -> str:
        file_ext = os.path.splitext(self.path)[1].lower()

        if file_ext == ".txt":
            self.text_file(self.path)
        elif file_ext == ".pdf":
            self.pdf_file(self.path)
        elif file_ext in (".doc", ".docx"):
            self.doc_file(self.path)
        elif file_ext in (".ppt", ".pptx"):
            self.ppt_file(self.path)

        return self.text_format.strip()

    def text_output(self) -> str:
        return self.convert()

    def text_file(self, path):
        try:
            with open(path, "r", encoding="utf-8") as text_file:
                self.text_format = text_file.read()
        except UnicodeDecodeError:
            with open(path, "r", encoding="latin-1") as text_file:
                self.text_format = text_file.read()

    def pdf_file(self, path):
        pdf_module = self._load_pdf_module()
        if pdf_module is None:
            raise RuntimeError("PDF extraction requires 'pymupdf'. Install dependencies from requirements.txt")

        doc = pdf_module.open(path)
        for page in doc:
            self.text_format += page.get_text("text")
            self.text_format += "\n"

    def doc_file(self, path):
        if docx is None:
            raise RuntimeError("Word extraction requires 'python-docx'. Install dependencies from requirements.txt")
        doc = docx.Document(path)
        for para in doc.paragraphs:
            if para.text:
                self.text_format += para.text + "\n"

    def ppt_file(self, path):
        if Presentation is None:
            raise RuntimeError("PowerPoint extraction requires 'python-pptx'. Install dependencies from requirements.txt")
        prs = Presentation(path)
        text_runs = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        self.text_format = "\n".join([value for value in text_runs if value])

    def _load_pdf_module(self):
        if pymupdf is not None:
            return pymupdf

        try:
            return importlib.import_module("pymupdf")
        except ImportError:
            try:
                return importlib.import_module("fitz")
            except ImportError:
                return None

          